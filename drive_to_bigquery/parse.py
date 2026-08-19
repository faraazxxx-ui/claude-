"""Turn downloaded bytes into DataFrames, and documents into text."""

from __future__ import annotations

import io
import json

import pandas as pd

from classify import sanitize_column_name

# Read everything as string first, then let BigQuery/pandas infer on the merged
# frame. Per-file inference is what causes schema drift across a family (one
# day's file has all-integer values, the next has a decimal).
_READ_KWARGS = {"dtype": str, "keep_default_na": False, "na_values": [""]}

MAX_TEXT_CHARS = 900_000  # keep a row comfortably under BigQuery's 100 MB cap


def normalize_columns(frame: pd.DataFrame) -> pd.DataFrame:
    """Sanitize headers and de-duplicate collisions."""
    seen: dict[str, int] = {}
    columns = []
    for position, raw in enumerate(frame.columns):
        name = sanitize_column_name(raw, position)
        if name in seen:
            seen[name] += 1
            name = f"{name}_{seen[name]}"
        else:
            seen[name] = 0
        columns.append(name)
    frame.columns = columns
    return frame


def read_tabular(data: bytes, fmt: str, name: str) -> list[tuple[str, pd.DataFrame]]:
    """Parse bytes into ``[(sheet_suffix, frame)]``.

    Multi-sheet workbooks yield one entry per sheet; everything else yields one
    entry with an empty suffix.
    """
    if fmt in {"csv", "tsv"}:
        sep = "\t" if fmt == "tsv" else ","
        frame = _read_csv(data, sep)
        return [("", normalize_columns(frame))]

    if fmt in {"xlsx", "xls", "xlsm"}:
        engine = "openpyxl" if fmt != "xls" else "xlrd"
        book = pd.read_excel(
            io.BytesIO(data), sheet_name=None, engine=engine, **_READ_KWARGS
        )
        out = []
        for sheet_name, frame in book.items():
            if frame.empty:
                continue
            out.append((str(sheet_name), normalize_columns(frame)))
        return out

    if fmt in {"json", "ndjson"}:
        return [("", normalize_columns(_read_json(data)))]

    raise ValueError(f"unsupported tabular format {fmt!r} for {name!r}")


def _read_csv(data: bytes, sep: str) -> pd.DataFrame:
    """Read a CSV, tolerating the encodings Takeout exports show up in."""
    last: Exception | None = None
    for encoding in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            return pd.read_csv(
                io.BytesIO(data),
                sep=sep,
                encoding=encoding,
                engine="python",
                on_bad_lines="skip",
                **_READ_KWARGS,
            )
        except (UnicodeDecodeError, pd.errors.ParserError) as exc:
            last = exc
    raise ValueError(f"could not parse CSV: {last}")


def _read_json(data: bytes) -> pd.DataFrame:
    """Flatten JSON into a frame, handling arrays, NDJSON, and wrapped objects."""
    text = data.decode("utf-8", errors="replace").strip()
    if not text:
        return pd.DataFrame()

    try:
        parsed = json.loads(text)
    except json.JSONDecodeError:
        # Assume NDJSON.
        rows = []
        for line in text.splitlines():
            line = line.strip()
            if not line:
                continue
            try:
                rows.append(json.loads(line))
            except json.JSONDecodeError:
                continue
        return pd.json_normalize(rows) if rows else pd.DataFrame()

    if isinstance(parsed, list):
        return pd.json_normalize(parsed)
    if isinstance(parsed, dict):
        # Takeout wraps the payload in a single key more often than not
        # (e.g. {"features": [...]} in Saved Places).
        list_values = [v for v in parsed.values() if isinstance(v, list)]
        if len(list_values) == 1:
            return pd.json_normalize(list_values[0])
        return pd.json_normalize([parsed])
    return pd.DataFrame({"value": [parsed]})


def extract_text(data: bytes, fmt: str) -> tuple[str, int, str]:
    """Return ``(text, page_count, method)`` for a document file."""
    if fmt in {"txt", "md", "html", "rtf"}:
        text = data.decode("utf-8", errors="replace")
        if fmt == "html":
            text = _strip_html(text)
        return text[:MAX_TEXT_CHARS], 0, f"decode:{fmt}"

    if fmt == "pdf":
        try:
            from pypdf import PdfReader
        except ImportError:
            return "", 0, "skipped:pypdf-missing"
        reader = PdfReader(io.BytesIO(data))
        pages = [(page.extract_text() or "") for page in reader.pages]
        return "\n".join(pages)[:MAX_TEXT_CHARS], len(pages), "pypdf"

    if fmt == "docx":
        try:
            import docx
        except ImportError:
            return "", 0, "skipped:python-docx-missing"
        document = docx.Document(io.BytesIO(data))
        text = "\n".join(p.text for p in document.paragraphs)
        return text[:MAX_TEXT_CHARS], 0, "python-docx"

    if fmt == "pptx":
        try:
            from pptx import Presentation
        except ImportError:
            return "", 0, "skipped:python-pptx-missing"
        deck = Presentation(io.BytesIO(data))
        chunks = []
        for slide in deck.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    chunks.append(shape.text_frame.text)
        return "\n".join(chunks)[:MAX_TEXT_CHARS], len(deck.slides), "python-pptx"

    return "", 0, f"unsupported:{fmt}"


def _strip_html(text: str) -> str:
    import re

    text = re.sub(r"<(script|style)[^>]*>.*?</\1>", " ", text, flags=re.S | re.I)
    text = re.sub(r"<[^>]+>", " ", text)
    return re.sub(r"\s{2,}", " ", text).strip()


def reconcile(frames: list[pd.DataFrame]) -> pd.DataFrame:
    """Union frames with differing columns into one, preserving every column.

    Missing columns become NULL rather than dropping the row — a day's file that
    lacks a column should not silently lose its other fields.
    """
    if not frames:
        return pd.DataFrame()
    ordered: list[str] = []
    for frame in frames:
        for column in frame.columns:
            if column not in ordered:
                ordered.append(column)
    aligned = [frame.reindex(columns=ordered) for frame in frames]
    return pd.concat(aligned, ignore_index=True)
